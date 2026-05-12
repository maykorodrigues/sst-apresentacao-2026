#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Gerador de Apresentação: Sucessso do Cliente - Maio 2026
Estrutura baseada na reunião com Lucas (12/05/2026)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Cores corporativas SST Card
COR_PRIMARIA = RGBColor(25, 135, 195)      # Azul
COR_SECUNDARIA = RGBColor(240, 80, 80)     # Vermelho/Alerta
COR_SUCESSO = RGBColor(76, 175, 80)        # Verde
COR_ATENCAO = RGBColor(255, 193, 7)        # Amarelo
COR_TEXTO = RGBColor(33, 33, 33)           # Cinza escuro
COR_FUNDO = RGBColor(245, 245, 245)        # Cinza claro

def adicionar_titulo_slide(slide, titulo, subtitulo=""):
    """Adiciona título e subtítulo ao slide"""
    title = slide.shapes.title
    title.text = titulo
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA

    if subtitulo and len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = subtitulo
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = COR_TEXTO

def adicionar_caixa_texto(slide, left, top, width, height, texto, tamanho=14, bold=False, cor=None):
    """Adiciona uma caixa de texto ao slide"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = texto
    p.font.size = Pt(tamanho)
    p.font.bold = bold
    if cor:
        p.font.color.rgb = cor
    else:
        p.font.color.rgb = COR_TEXTO
    return txBox

# Criar apresentação
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ============ SLIDE 1: CAPA ============
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COR_PRIMARIA

# Título
adicionar_caixa_texto(slide1, Inches(1), Inches(2.5), Inches(8), Inches(1.5),
    "SUCESSO DO CLIENTE", tamanho=60, bold=True, cor=RGBColor(255, 255, 255))

# Subtítulo
adicionar_caixa_texto(slide1, Inches(1), Inches(4), Inches(8), Inches(1),
    "SST Card — Bairro da Paz | Maio 2026", tamanho=32, cor=RGBColor(255, 255, 255))

# Rodapé
adicionar_caixa_texto(slide1, Inches(1), Inches(6.5), Inches(8), Inches(0.8),
    "Apresentação para Rogério Ferreira | Mayko Rodrigues", tamanho=14, cor=RGBColor(200, 200, 200))

# ============ SLIDE 2: PROBLEMA CENTRAL ============
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
adicionar_titulo_slide(slide2, "🎯 O Problema Central", "Churn Descontrolado")

adicionar_caixa_texto(slide2, Inches(0.5), Inches(1.8), Inches(9), Inches(0.6),
    "Temos 27 cancelamentos documentados, mas:", tamanho=16, bold=True)

# Bullet points
bullets = [
    "❌ Motivos não mapeados — por que cancelaram?",
    "❌ Processo de retenção inexistente",
    "❌ Reclamações dispersas entre recepção, Raquel, WhatsApp",
    "❌ Sem dono formal de 'Customer Success'",
    "❌ Possível churn oculto entre 1.293 inativos"
]

y_pos = 2.6
for bullet in bullets:
    adicionar_caixa_texto(slide2, Inches(1), Inches(y_pos), Inches(8.5), Inches(0.4),
        bullet, tamanho=14)
    y_pos += 0.5

adicionar_caixa_texto(slide2, Inches(0.5), Inches(6.8), Inches(9), Inches(0.6),
    "⚠️ Se não mapearmos agora, não sabemos o que ajustar.", tamanho=12, bold=True, cor=COR_SECUNDARIA)

# ============ SLIDE 3: NÚMEROS DO SISTEMA ============
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
adicionar_titulo_slide(slide3, "📊 Cenário Atual", "Contratos no Boom")

# Tabela com dados
dados = [
    ("Total de Contratos", "1.708", COR_TEXTO),
    ("Ativos (Pagando)", "415", COR_SUCESSO),
    ("Inativos/Indefinidos", "1.293", COR_ATENCAO),
    ("Cancelados (COM TAG)", "27", COR_SECUNDARIA),
]

y_start = 1.8
for label, valor, cor in dados:
    adicionar_caixa_texto(slide3, Inches(1), Inches(y_start), Inches(4), Inches(0.4),
        label, tamanho=14, bold=True)
    adicionar_caixa_texto(slide3, Inches(5.5), Inches(y_start), Inches(3), Inches(0.4),
        valor, tamanho=20, bold=True, cor=cor)
    y_start += 0.7

# Churn calculado
adicionar_caixa_texto(slide3, Inches(1), Inches(4.8), Inches(9), Inches(0.5),
    "Churn Rate Atual: 27 ÷ 1.708 = 1,58% ✅ (Saudável)", tamanho=14, bold=True, cor=COR_SUCESSO)

adicionar_caixa_texto(slide3, Inches(1), Inches(5.6), Inches(9), Inches(1),
    "⚠️ Mas: 1.293 inativos = Churn potencial desconhecido. Precisamos mapear quem são, por quê estão inativos, e quantos realmente cancelaram.",
    tamanho=12, cor=COR_ATENCAO)

# ============ SLIDE 4: OS 27 CANCELADOS (EXEMPLO) ============
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
adicionar_titulo_slide(slide4, "🔍 Exemplo: Eliane", "Primeira análise de cancelamento")

exemplo_dados = [
    ("Nome", "Eliane"),
    ("CPF", "***.***.***-**"),
    ("Tempo de cliente", "10 meses"),
    ("Motivo", "Mudança de cidade"),
    ("Data cancelamento", "19/03/2026"),
    ("Canal", "Chat/Recepção + Karine"),
    ("Multa cobrada", "R$ 49,80 (2 meses restantes)"),
    ("Reversível?", "NÃO (mudança permanente)"),
    ("Dependentes", "Sim (mãe, 88 anos)"),
    ("Histórico clínica", "1 oftalmo em 2024"),
]

y_start = 1.8
for label, valor in exemplo_dados:
    adicionar_caixa_texto(slide4, Inches(1), Inches(y_start), Inches(3.5), Inches(0.3),
        label + ":", tamanho=11, bold=True)
    adicionar_caixa_texto(slide4, Inches(4.7), Inches(y_start), Inches(4.5), Inches(0.3),
        valor, tamanho=11)
    y_start += 0.3

adicionar_caixa_texto(slide4, Inches(1), Inches(6.5), Inches(9), Inches(0.8),
    "💡 Insight: Temos 26 mais para mapear. Eliane é legítima, mas outras podem ser reversíveis (não usa benefício, problemas financeiros, etc.).",
    tamanho=11, cor=COR_PRIMARIA)

# ============ SLIDE 5: SOLUÇÃO — PLANILHA CHURNING ============
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
adicionar_titulo_slide(slide5, "📋 Solução 1: Planilha Churning", "Controle centralizado")

colunas = [
    "CPF",
    "Nome",
    "Data Cancelamento",
    "MOTIVO ⭐",
    "Canal",
    "Reversível?",
    "Status",
]

y_start = 1.8
adicionar_caixa_texto(slide5, Inches(0.5), Inches(y_start), Inches(9), Inches(0.4),
    " | ".join(colunas), tamanho=10, bold=True, cor=COR_PRIMARIA)
y_start += 0.5

motivos_exemplo = [
    "Mudança de cidade",
    "Não usa benefício",
    "Problemas financeiros",
    "Outro plano de saúde",
    "Mal atendimento (clínica)",
]

for motivo in motivos_exemplo:
    adicionar_caixa_texto(slide5, Inches(0.5), Inches(y_start), Inches(9), Inches(0.3),
        f"... | ... | ... | {motivo} | ... | ... | ...", tamanho=9)
    y_start += 0.4

adicionar_caixa_texto(slide5, Inches(1), Inches(5.5), Inches(8), Inches(1.5),
    "🎯 OBJETIVO:\n1. Exportar 27 do Boom\n2. Criar Google Sheet + Notion\n3. Lucas ligar para cada um\n4. Documentar TODOS os motivos",
    tamanho=12, bold=True, cor=COR_PRIMARIA)

# ============ SLIDE 6: PROCESSO DE RETENÇÃO ============
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
adicionar_titulo_slide(slide6, "☎️ Solução 2: Processo de Retenção", "Estilo 'Cancelamento de Telefonia'")

processo_steps = [
    "1️⃣ Cliente chega na recepção: 'Quero cancelar'",
    "2️⃣ Recepcionista: 'Lucas vai ligar em 24-48h'",
    "3️⃣ Lucas liga (WhatsApp Voice ou chip)",
    "4️⃣ Conversa: Motivo + Validação + Oferta",
    "5️⃣ Resultado: Reversão ✅ ou Saída Documentada ✅",
]

y_start = 1.8
for step in processo_steps:
    adicionar_caixa_texto(slide6, Inches(1.5), Inches(y_start), Inches(7), Inches(0.4),
        step, tamanho=14, bold=True)
    y_start += 0.7

adicionar_caixa_texto(slide6, Inches(1), Inches(5.3), Inches(8), Inches(1.8),
    "OFERTAS DE RETENÇÃO:\n• Voucher de consulta\n• Explicação de benefícios não usados (telemedicina, auxílio federal)\n• Negociação financeira\n• Agenda consulta de follow-up",
    tamanho=11, cor=COR_SUCESSO)

# ============ SLIDE 7: RESPONSABILIDADES LUCAS ============
slide7 = prs.slides.add_slide(prs.slide_layouts[1])
adicionar_titulo_slide(slide7, "👤 Funções Formalizadas de Lucas", "Customer Success Owner")

funcoes = [
    ("Desfiliação (Cancelamento)", "Ligar para os 27, entender motivos", "Churn < 5%"),
    ("Reclamações", "Receber de Raquel, processar, informar Rogério", "% resolvidas"),
    ("Onboarding", "Ligar ao novo cliente (script boas-vindas)", "100% ligados"),
    ("Retenção", "Se vai cancelar, tentar reverter", "≥30% revertidos"),
]

y_start = 1.8
for funcao, acao, kpi in funcoes:
    adicionar_caixa_texto(slide7, Inches(0.5), Inches(y_start), Inches(2.8), Inches(0.6),
        funcao, tamanho=11, bold=True, cor=COR_PRIMARIA)
    adicionar_caixa_texto(slide7, Inches(3.5), Inches(y_start), Inches(3.2), Inches(0.6),
        acao, tamanho=10)
    adicionar_caixa_texto(slide7, Inches(7), Inches(y_start), Inches(2.5), Inches(0.6),
        kpi, tamanho=10, bold=True, cor=COR_SUCESSO)
    y_start += 1

# ============ SLIDE 8: MÉTRICAS DE LUCAS ============
slide8 = prs.slides.add_slide(prs.slide_layouts[1])
adicionar_titulo_slide(slide8, "📈 KPIs de Customer Success", "Como Lucas será medido")

metricas = [
    ("1. Filiações Processadas", "100% ligados em D+1"),
    ("2. Desfiliações Atendidas", "100% dos que chamam"),
    ("3. Churn Rate", "< 5% (benchmark: < 3%)"),
    ("4. Reversões (Retenção)", "≥ 30% dos que cancelam"),
    ("5. NPS (Satisfação)", "≥ 8 (promotores)"),
]

y_start = 1.8
for metrica, kpi in metricas:
    adicionar_caixa_texto(slide8, Inches(1), Inches(y_start), Inches(4.5), Inches(0.4),
        metrica, tamanho=12, bold=True, cor=COR_PRIMARIA)
    adicionar_caixa_texto(slide8, Inches(5.8), Inches(y_start), Inches(3.5), Inches(0.4),
        kpi, tamanho=12, cor=COR_SUCESSO)
    y_start += 0.8

# ============ SLIDE 9: PRÓXIMOS PASSOS (24-72h) ============
slide9 = prs.slides.add_slide(prs.slide_layouts[1])
adicionar_titulo_slide(slide9, "🚀 Próximos Passos", "Execução imediata")

passos = [
    ("HOJE 14h", "Exportar 27 cancelados do Boom", "Lucas", "Planilha em mãos"),
    ("HOJE 15h", "Criar Planilha 'Churning' Google Sheets", "Mayko", "Link compartilhado"),
    ("Amanhã", "Sincronizar com Notion", "Mayko", "Dashboard em Notion"),
    ("Esta semana", "Lucas ligar para 5 dos 27", "Lucas", "Motivos documentados"),
    ("Próx. reunião", "Trazer 27 motivos mapeados", "Lucas", "Tabela completa"),
]

y_start = 1.8
for data, acao, responsavel, metrica in passos:
    adicionar_caixa_texto(slide9, Inches(0.5), Inches(y_start), Inches(1.8), Inches(0.4),
        data, tamanho=10, bold=True, cor=COR_PRIMARIA)
    adicionar_caixa_texto(slide9, Inches(2.5), Inches(y_start), Inches(3.5), Inches(0.4),
        acao, tamanho=10)
    adicionar_caixa_texto(slide9, Inches(6.2), Inches(y_start), Inches(1.5), Inches(0.4),
        responsavel, tamanho=9, cor=COR_ATENCAO)
    adicionar_caixa_texto(slide9, Inches(7.9), Inches(y_start), Inches(1.7), Inches(0.4),
        metrica, tamanho=9, cor=COR_SUCESSO)
    y_start += 0.8

# ============ SLIDE 10: DECISÃO ============
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide10.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COR_SUCESSO

adicionar_caixa_texto(slide10, Inches(1), Inches(2), Inches(8), Inches(1.5),
    "✅ LUCAS COMO DONO FORMAL\nDE CUSTOMER SUCCESS?", tamanho=48, bold=True,
    cor=RGBColor(255, 255, 255))

adicionar_caixa_texto(slide10, Inches(1), Inches(3.8), Inches(8), Inches(2),
    "✓ Desfiliação + Reclamações + NPS\n✓ Responsável pelo Churn Rate\n✓ Base para supervisionário em Bairro da Paz",
    tamanho=18, cor=RGBColor(255, 255, 255))

adicionar_caixa_texto(slide10, Inches(1), Inches(6), Inches(8), Inches(1),
    "Aprovado por: ___________________     Data: ___________",
    tamanho=14, cor=RGBColor(255, 255, 255))

# ============ SLIDE 11: OBRIGADO ============
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide11.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COR_PRIMARIA

adicionar_caixa_texto(slide11, Inches(1), Inches(2.5), Inches(8), Inches(2),
    "OBRIGADO", tamanho=72, bold=True, cor=RGBColor(255, 255, 255))

adicionar_caixa_texto(slide11, Inches(1), Inches(4.5), Inches(8), Inches(1.5),
    "Mayko Rodrigues\nConsultor RevOps — Sucesso de Todos\ncontato@sucessodetodos.com.br",
    tamanho=16, cor=RGBColor(255, 255, 255))

# Salvar
output_path = r"C:\Users\mayko\meu-cerebro\01-projetos\consultoria-comercial\clientes\SST_Clinica_Bairro_da_Paz\lucas-cs\Apresentação-Sucesso-do-Cliente-Maio-2026-v2.pptx"
prs.save(output_path)

print(f"[OK] Apresentacao salva em: {output_path}")
print("[OK] 11 slides criados com dados contextuais da reuniao com Lucas")
